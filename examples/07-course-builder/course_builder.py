"""
Example 07: Course Builder — End-to-end course construction on Feishu Wiki.

This example demonstrates a complete workflow:
  1. Create course outline (chapter + section nodes)
  2. Fill content into each section
  3. Generate PPT slides
  4. Upload PPTs to Drive and move into Wiki

Usage:
    export FEISHU_APP_ID="cli_xxx"
    export FEISHU_APP_SECRET="xxx"
    export WIKU_SPACE_ID="your_space_id"
    python course_builder.py

Dependencies:
    pip install feishu-kit python-pptx
"""

import asyncio
import os
import tempfile
from dataclasses import dataclass

from feishu_kit import FeishuClient
from feishu_kit.modules.drive import DriveService
from feishu_kit.modules.wiki import WikiService

# ── Block helpers ──────────────────────────────────────────


def t(content: str, bold: bool = False, italic: bool = False) -> dict:
    style = {}
    if bold:
        style["bold"] = True
    if italic:
        style["italic"] = True
    return {"text_run": {"content": content, "text_element_style": style}}


def heading2(text: str) -> dict:
    return {"block_type": 4, "heading2": {"elements": [t(text)], "style": {}}}


def heading3(text: str) -> dict:
    return {"block_type": 5, "heading3": {"elements": [t(text)], "style": {}}}


def paragraph(*parts) -> dict:
    return {"block_type": 2, "text": {"elements": list(parts), "style": {}}}


def bullet(*parts) -> dict:
    return {"block_type": 12, "bullet": {"elements": list(parts), "style": {}}}


# ── PPT helper ─────────────────────────────────────────────


def make_ppt(title: str, slides: list[tuple[str, list[str]]], path: str):
    """Create a simple PPT file. slides = [(slide_title, [bullet_points])]."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True

    for slide_title, bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Slide title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(32)
        p.font.bold = True
        # Bullets
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
        tf2 = txBox2.text_frame
        for i, bullet_text in enumerate(bullets):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = bullet_text
            p.font.size = Pt(22)
            p.space_after = Pt(10)

    prs.save(path)


# ── Course definition ──────────────────────────────────────


@dataclass
class Section:
    title: str
    content_blocks: list[dict]
    ppt_slides: list[tuple[str, list[str]]]


@dataclass
class Chapter:
    title: str
    sections: list[Section]


# Demo course outline
COURSE = [
    Chapter(
        "绪论",
        [
            Section(
                "1.1 什么是具身智能",
                content_blocks=[
                    heading2("具身智能的概念"),
                    paragraph(
                        t("具身智能（Embodied AI）", bold=True),
                        t("是指智能体通过与物理环境的交互来获取感知、执行行动的智能范式。"),
                    ),
                    paragraph(
                        "与传统「离身」AI 不同，具身智能要求智能体拥有身体，能通过传感器感知环境，通过执行器改变环境。"
                    ),
                    heading3("核心特征"),
                    bullet(
                        t("感知—行动闭环", bold=True), t("：通过行动改变环境，再感知变化，形成闭环")
                    ),
                    bullet(t("环境交互性", bold=True), t("：智能行为在真实或仿真物理环境中实现")),
                    bullet(t("多模态感知", bold=True), t("：融合视觉、听觉、触觉等多种传感器信息")),
                    bullet(t("实时性与鲁棒性", bold=True), t("：在不确定条件下保持稳定运行")),
                    bullet(t("学习与适应", bold=True), t("：通过持续交互学习新技能、适应新场景")),
                    heading3("生活中的类比"),
                    paragraph(
                        t("会开车 vs 懂交规：", bold=True),
                        t(
                            "熟读驾驶手册但从未上路的人，和有实际驾驶经验的人，差距在哪？后者在与环境的交互中积累了「身体智慧」—— 这正是具身智能的核心。"
                        ),
                    ),
                    heading2("参考文献"),
                    bullet(
                        t("Brooks, R. (1991). Intelligence Without Representation. "),
                        t("Artificial Intelligence", italic=True),
                        t(", 47, 139-160."),
                    ),
                ],
                ppt_slides=[
                    (
                        "什么是具身智能",
                        [
                            "Embodied AI = 智能体 + 物理环境 + 持续交互",
                            "与传统AI的核心区别：拥有身体，与环境交互",
                            "感知→认知→决策→执行→感知 闭环",
                        ],
                    ),
                    (
                        "核心特征",
                        [
                            "感知—行动闭环：主动获取信息而非被动接收",
                            "环境交互性：在真实物理环境中实现智能",
                            "多模态感知：视觉+听觉+触觉+力觉融合",
                            "实时性与鲁棒性：应对噪声和动态变化",
                            "学习与适应：从经验中持续提升",
                        ],
                    ),
                    (
                        "参考文献",
                        [
                            "Brooks (1991) - Intelligence Without Representation",
                            "Pfeifer & Bongard (2007) - How the Body Shapes the Way We Think",
                        ],
                    ),
                ],
            ),
            Section(
                "1.2 发展历程",
                content_blocks=[
                    heading2("具身智能的发展历程"),
                    bullet(
                        t("1948年", bold=True),
                        t("：Wiener 出版《Cybernetics》，控制论奠定感知-行动闭环理论基础"),
                    ),
                    bullet(
                        t("1961年", bold=True), t("：Unimate 安装于 GM 工厂，全球第一台工业机器人")
                    ),
                    bullet(t("1991年", bold=True), t("：Brooks 提出「世界本身就是最好的模型」")),
                    bullet(
                        t("2016年", bold=True), t("：AlphaGo 展示从交互中超越人类先验知识的可能性")
                    ),
                    bullet(
                        t("2023年", bold=True),
                        t("：Google RT-2 将大模型与机器人控制统一，VLA 模型诞生"),
                    ),
                    heading2("参考文献"),
                    bullet(t("Wiener, N. (1948). Cybernetics. MIT Press.")),
                    bullet(
                        t(
                            "Silver, D. et al. (2016). Mastering the Game of Go. Nature, 529, 484-489."
                        )
                    ),
                ],
                ppt_slides=[
                    (
                        "发展历程",
                        [
                            "1948: 控制论 (Wiener)",
                            "1961: 第一台工业机器人 (Unimate)",
                            "1991: 包容架构 (Brooks)",
                            "2016: AlphaGo (DeepMind)",
                            "2023: RT-2 / VLA 模型 (Google)",
                        ],
                    ),
                ],
            ),
        ],
    ),
]


# ── Main workflow ──────────────────────────────────────────


async def main():
    space_id = os.environ.get("WIKU_SPACE_ID", "")

    async with FeishuClient(
        os.environ["FEISHU_APP_ID"],
        os.environ["FEISHU_APP_SECRET"],
    ) as client:
        wiki = WikiService(client)
        drive = DriveService(client)

        for chapter in COURSE:
            print(f"\n{'=' * 50}")
            print(f"章节: {chapter.title}")
            print(f"{'=' * 50}")

            # Step 1: Create chapter node
            result = await wiki.create_node(space_id, title=chapter.title)
            parent_token = result["data"]["node"]["node_token"]
            print(f"  创建节点: {parent_token}")

            for section in chapter.sections:
                print(f"\n  --- {section.title} ---")

                # Step 2: Create section node
                result = await wiki.create_node(
                    space_id,
                    title=section.title,
                    parent_node_token=parent_token,
                )
                obj_token = result["data"]["node"]["obj_token"]
                node_token = result["data"]["node"]["node_token"]
                print(f"  创建文档: obj={obj_token}")

                # Step 3: Fill content
                await wiki.create_doc_block(obj_token, obj_token, section.content_blocks)
                print(f"  写入内容: {len(section.content_blocks)} blocks")

                # Step 4: Create PPT
                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
                    ppt_path = f.name
                make_ppt(section.title, section.ppt_slides, ppt_path)
                print(f"  生成PPT: {ppt_path}")

                # Step 5: Upload PPT to Drive
                with open(ppt_path, "rb") as f:
                    file_data = f.read()
                ppt_name = f"{section.title}.pptx"
                upload_result = await drive.upload_file("", ppt_name, file_data)
                file_token = upload_result.get("data", {}).get("file_token", "")
                print(f"  上传云盘: {file_token}")

                # Step 6: Move PPT into Wiki
                if file_token:
                    move_result = await wiki.move_docs_to_wiki(
                        space_id=space_id,
                        parent_wiki_token=node_token,
                        obj_token=file_token,
                        obj_type="file",
                    )
                    print(f"  移入Wiki: code={move_result.get('code')}")

                # Cleanup temp file
                os.unlink(ppt_path)

        print(f"\n{'=' * 50}")
        print("课程构建完成!")
        print(f"{'=' * 50}")


if __name__ == "__main__":
    asyncio.run(main())
